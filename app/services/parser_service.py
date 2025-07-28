# backend/app/services/parser_service.py
import os
import fitz  # PyMuPDF
from docx import Document
import pptx
import logging

class AppError(Exception):
    pass

def parse_file(path: str) -> str:
    """Parses the content of a file based on its extension."""
    try:
        extension = os.path.splitext(path)[1].lower()
        if extension in ['.txt', '.md', '.rs', '.js', '.ts', '.py', '.html', '.css', '.json', '.toml']:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
        elif extension == '.pdf':
            with fitz.open(path) as doc:
                text = ""
                for page in doc:
                    text += page.get_text()
                return text
        elif extension == '.docx':
            doc = Document(path)
            return "\n".join(para.text for para in doc.paragraphs)
        elif extension == '.pptx':
            prs = pptx.Presentation(path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            return "\n".join(text_runs)
        else:
            logging.warning(f"Unsupported file type for parsing: {extension}")
            raise AppError(f"Unsupported file type: {extension}")
    except Exception as e:
        logging.error(f"Failed to parse file {path}: {e}")
        raise AppError(f"Failed to parse file {path}: {e}")